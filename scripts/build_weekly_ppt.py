"""Build weekly work report PPT from the 2026-04-14~2026-04-20 Obsidian markdown."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(r"docs/weekly_reports/주간업무보고_2026-04-14_2026-04-20.pptx")

# Design tokens
ACCENT = RGBColor(0x1F, 0x4E, 0x79)      # deep navy
ACCENT2 = RGBColor(0x2E, 0x75, 0xB6)     # azure
MUTED = RGBColor(0x59, 0x59, 0x59)
TEXT = RGBColor(0x1F, 0x1F, 0x1F)
BG_SOFT = RGBColor(0xF2, 0xF2, 0xF2)
DANGER = RGBColor(0xC0, 0x50, 0x4D)
OK = RGBColor(0x2E, 0x8B, 0x57)
BORDER = RGBColor(0xBF, 0xBF, 0xBF)

FONT_KO = "맑은 고딕"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def set_run(run, *, size=None, bold=None, color=None, font=FONT_KO):
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = font
    # Set east asian font
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            from lxml import etree
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", font)


def add_textbox(slide, left, top, width, height, text, *, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        set_run(r, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, left, top, width, height, *, fill=ACCENT, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_header(slide, title, subtitle=None):
    # accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=ACCENT)
    # title
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12.5), Inches(0.6),
                title, size=26, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.95), Inches(12.5), Inches(0.35),
                    subtitle, size=13, color=MUTED)
    # divider
    add_rect(slide, Inches(0.6), Inches(1.35), Inches(12.1), Emu(12700),
             fill=ACCENT2)


def new_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank)


def add_footer(slide, page, total):
    add_textbox(slide, Inches(0.6), Inches(7.1), Inches(6), Inches(0.3),
                "연구개발팀 · 백현규 · 2026-04-14 ~ 2026-04-20", size=9, color=MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.3), Inches(0.3),
                f"{page} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ------------------------------------------------------------
# Content
# ------------------------------------------------------------
TITLE = "주간업무보고"
PERIOD = "2026년 4월 3주차 (04/14 화 ~ 04/20 월)"
AUTHOR = "연구개발팀 · 백현규"

EXEC_SUMMARY = [
    "소부장 R&D 사업계획서 본문 8개 섹션 신규 작성 완료 (과제 개요·연구방법·WBS·활용방안·실증계획·지역특화·개방형혁신·선행연구)",
    "ATL AXIS 일정 관리 풀스택 고도화 완료 (실제 시작/종료일 필드·간트 이중 바·긴급도 6등급 시각화·말줄임 툴팁)",
    "ATL-RAMS ALE·ALDC·ALC 파형 데이터 손실·오실로스코프·복합파형·연결 flicker·반응형 레이아웃 등 다수 품질 결함 일괄 수정 (996 테스트 통과)",
    "ATL AXIS 배포 환경 도메인 연결 (nginx + Cloudflare TLS), RBAC 풀스택, n8n 자동화 인프라 (19개 이벤트·워크플로우 4종) 구축",
    "다음 주 포커스: 시험 전/후 속성 QA 안정화, ALE 복합파형 기능 확장, lineauto PR-8 후속 로드맵",
]

PROJECTS = [
    ("소부장 R&D 사업계획서", "ATL-RAMS 확장 연구개발 과제", [
        ("사업계획서 목차 리스트업 및 WBS 수립", "04/14", "100%",
            "공고문·제출 서류·FAQ 등 참고자료 7건 분석 기반 본문 전체 목차 구조화. ATL-RAMS 현재 개발 현황(완성도 70~95%) 및 W1~W6 마일스톤과의 연계 관계 명시"),
        ("선행연구개발 결과 활용계획 신규 작성", "04/14", "100%",
            "HW/SW 미개발 전제로 '고속 파형 생성 기술 선행연구로 원천기술 확보, 본 과제에서 HW+SW 신규 개발' 컨셉 확립. 6개 원천기술×3단 매핑 테이블, TRL 3~4 조정"),
        ("연구개발 방법 섹션 신규 작성 및 4대 구성요소 기반 재편", "04/15", "100%",
            "4대 시스템 구성요소(ATL_RAMS·WaveGenerator·AMP·챔버) 중심 재구성. 챔버 영역을 '전문 제조사 상용 열챔버 도입 + 통신 인터페이스·SW 연동' 컨셉으로 전환"),
        ("기술개발일정 및 WBS 작성", "04/15", "100%",
            "7개 Work Package 구조, M1~M24 월별 Gantt Chart, 5개 핵심 마일스톤 정의. 챔버(WP3) 일정 2차년도 이동 및 의존 관계 재조정"),
        ("연구개발성과 활용방안 및 기대효과 작성", "04/15", "100%",
            "사업화 목표·전략, 지정학적 리스크 대응(호르무즈·사우디 Vision 2030), ESG 연계, IP 전략, 기대효과 구성. 제품 라인업을 ATL-RAMS ALE 단일 제품으로 통합"),
        ("지역특화·실증계획·개방형혁신·과제 개요 문서 작성", "04/15", "100%",
            "전기전자 분야 전략 품목 정의 및 1:1 대응표. 4단계 실증 로드맵 수립(내부→통합→현장→공인). 타 분야 확장 대응표(자동차 EMC·의료기기·배터리 SOH·우주항공)"),
    ]),
    ("ATL AXIS (Frontend + Backend)", "SaaS 시험소 운영 플랫폼", [
        ("일정 관리 실제 시작/종료일 필드 및 간트 이중 바 구현", "04/14", "100%",
            "백엔드 엔티티·DTO·Service 확장, 프론트엔드 UI 5종(다이얼로그·상세 시트·간트·목록·Excel) 확장. 간트 이중 바(실제 solid + 계획 dashed) + 툴팁 병렬 표시"),
        ("배포 환경 도메인 연결 및 nginx 리버스 프록시 적용", "04/16", "100%",
            "atl-axis.com 도메인 80/443 포트 리스너 부재 문제 해결. Cloudflare Origin Certificate + TLS 1.2/1.3, Next.js 정책 기반 재작성(캐싱·보안 헤더·gzip·CORS)"),
        ("내 일정 관리 품질 결함 일괄 수정", "04/16~17", "100%",
            "Radix Collapsible onOpenChange 패턴 수정, 간트 트리 전체 접기/펼치기 복구, Hydration Mismatch 수정, TruncatedText 말줄임+호버 툴팁"),
        ("긴급도(Urgency) 시각화 구현", "04/17", "100%",
            "6등급 체계(overdue/due-today/at-risk/in-progress/upcoming/completed) 통합 판정 헬퍼 구축. 상단 공통 배너·재사용 뱃지·간트/리스트/보드/캘린더 4개 뷰 일괄 적용"),
        ("RBAC 풀스택 구현 + 규격·시험항목·견적 3개 페이지 품질 점검", "04/17", "100%",
            "tb_user_role 시드 마이그레이션, AuthService 권한 합산, 14개 컨트롤러 @Roles(). Next.js middleware·RouteGuard·PermissionGate·/403 페이지 구축"),
        ("n8n 자동화 공통 웹훅 인프라 + Phase 2~8 확장 + 워크플로우 4종", "04/17", "100%",
            "17개 이벤트·HMAC-SHA256·5회 재시도 크론, n8n-auth.guard. 마스터 라우터·Slack Interactive·DocuSign·Google Calendar Sync 워크플로우"),
        ("캘린더 API 500·LEG 편집 500·규격 삭제 400 에러 수정", "04/16~17", "100%",
            "actual_start_date_time 누락 컬럼 마이그레이션, uq_leg_code UNIQUE 슬롯 해제 패턴. DeleteReasonDialog 조기 종료 해소 및 HAS_REFERENCES code propagation"),
        ("라우팅 세션 동기화 버그 수정", "04/17", "100%",
            "App Router route group((authenticated)) 도입으로 children remount 경쟁 원천 제거"),
        ("알림 기능 풀스택 구현", "04/17", "100%",
            "NotificationsGateway·EventEmitterModule·ScheduleModule 리마인더 크론. REST·Zustand+Immer·WebSocket 훅·NotificationBell 구축"),
        ("시험 항목 시험 전/후(Phase) 속성 풀스택 구현 및 QA 통과", "04/20", "100%",
            "DB 컬럼·enum·DTO·서비스 확장, TestPhaseBadge·RadioGroup·LegOverviewGrid 3분할 리팩터. Phase(1차)→ItemRole(2차) 2단 그룹핑, LEG 빌더 Phase×Zone AND 필터"),
        ("디자인 시스템 시맨틱 컬러 토큰 전역 도입", "04/20", "100%",
            "24+ 파일의 bg-emerald/amber/rose/sky/indigo/violet/teal/cyan 패턴 해소. --color-success/warning/info + foreground 6종 매핑 신설, 라이트/다크 자동 분기"),
        ("견적 빌더·이메일·PDF 기능 개선", "04/20", "100%",
            "LEG zone 필터 구현 및 이후 실내/실외 배지·Zone 필터 제거 리팩터. 네비게이터 접기·오버뷰 확대, 이메일 금액 제거, PDF 조항 번호 표시, CC 자동완성 구현"),
        ("규격·사용자·계정 관련 버그 4종 수정", "04/20", "100%",
            "소프트삭제 UNIQUE 충돌 → ConflictException(409) 변환. SUPER_ADMIN 403 차단 해소, NavUser hydration mismatch 수정, 규격 빌더 pseudo 상태 수정"),
        ("DB 수정 후 list 화면 stale UI 일괄 해결", "04/20", "100%",
            "useRefetchOnFocus 공통 훅 신설(mount + visibilitychange + focus, 1s throttle). 6개 list 페이지 일괄 교체 + create store 강제 동기화"),
    ]),
    ("ATL-RAMS (JavaFX + Spring Boot)", "EMC 시험 파형 생성 장비", [
        ("파형·실시간 모니터링 데이터 손실 문제 해결", "04/16", "100%",
            "ALE/ALC/ALDC 이중 버퍼(DisplayWindow/RetentionStore) 아키텍처 적용. ALC O(n²) 루프 제거, ALDC 채널 토글 정책 전환, ALE LIVE Freeze/Pan. 891 테스트 통과"),
        ("ALE 모듈 QA 테스트 인프라 구축 및 커버리지 보강", "04/17", "100%",
            "JaCoCo·TestFX·Testcontainers·Awaitility 도입. Strategy 10종 계약 테스트 40케이스, Rule 헬퍼 3종. 996 tests 전부 GREEN"),
        ("ALE 가상 오실로스코프 반응형·기능 결함 일괄 수정", "04/17", "100%",
            "FXML 반응형 전환, ScreenSizeUtil 적용. 자동 TIME/DIV 락 피드백 루프, LIVE 모드 무반응, Seek 드래그 덮어쓰기 방지. 274 tests 0 failures"),
        ("ALE 장비 연결 상태 UI 깜빡임 버그 수정", "04/17", "100%",
            "폴러 양방향 전환을 positive 승격만으로 단순화, HEARTBEAT_TIMEOUT_MS 10→15s 완화"),
        ("ALE 다크 모드 UI/UX 고도화", "04/17", "100%",
            "4단계 엘리베이션 토큰 체계, 10개 패널 재매핑. WCAG AA 기준(4.5:1) 충족"),
        ("JavaFX 반응형 레이아웃 전면 수정 및 전체화면 복원", "04/17", "100%",
            "ScreenSizeUtil 신규, 4개 Controller + 8개 FXML 하드코딩 제거. 25개 다이얼로그 setMaxWidth/Height, maximized 복원 로직"),
        ("ALE 복합파형 미리보기 LINEAR/EXR 보정 및 구간 반복", "04/20", "100%",
            "LINEAR 수직 강제 절대 10ms OR 비율 2%로 확장, EXR 최소 표시 폭 보장. innerRepeat 풀스택(Flyway V44/V46·FXML UI·검증 불변식·3테마 CSS)"),
        ("ALE 파형 미리보기 전체 영역 채우기", "04/16", "100%",
            "NumberAxis autoRanging 한계 해소(명시적 bound + tickUnit 수동 지정). resetZoom 호환"),
    ]),
    ("ATL-RAMS RFQ 문서 (ALDC/ALPV)", "글로벌 벤더 견적 자료", [
        ("ALDC·ALPV RFQ 비교 견적 자료 산출 및 재설계", "04/17", "100%",
            "Keysight·HIOKI·Chroma·Kikusui·Yokogawa·Ametek 6개 벤더 비교. 4개 카테고리 구분, 조합 3안 비교 → 조합 B 권장. Markdown/Excel 6종 산출"),
        ("ALDC 채널당·ALPV 유닛당 원가 구조 재설계", "04/17", "100%",
            "ALDC ₩158,500 산출(원가율 5~8%). ALPV ₩1,195,600 산출(원가율 12%). R&D 투자 회수 분석 추가"),
        ("ALPV RFQ 조합 B R&D 투자 규모 재조정", "04/17", "100%",
            "판매가 ₩10M 대비 조합 B ₩64.5M이 6.4배 과중 지적 반영. Cat A·C·D 다운스케일, 조합 B 총액 ₩64.5M→₩41.4M"),
        ("ALDC RFQ 유닛당 ₩10M 기준 전면 재설계", "04/17", "100%",
            "채널당→유닛당 가격 전략 통일, 매출총이익률 95.2%. ALPV_RFQ.xlsx 요약 시트 합계·VAT·권장 문구 갱신"),
    ]),
    ("lineauto (리니지 클래식 자동사냥 매크로)", "Python 자동화 봇", [
        ("Phase 1 자동공격 MVP (Click+Drag-Down+Release)", "04/16", "100%",
            "HIDDevice Protocol drag 메서드 추가, 5개 백엔드 일괄 구현. DragTimingGenerator·BT drag provider·Engine. 564 passed"),
        ("실전 로그 기반 버그 수정 + Phase 2 HuntingLoop + Phase 3 APPROACHING/LOOTING", "04/17", "100%",
            "클릭 스팸 제거, Drag-once + click muting. vision/tracker.py 신규(IoU+class match). HuntingLoop 5-state FSM 완성. 634 passed"),
        ("OCR 툴팁 hover-probe 기반 몬스터 이름 확정 타겟팅", "04/17", "100%",
            "NameProbe.probe API(250ms timeout·cache_ttl_ms=1500). NearestTargeting·is_reachable 직선 샘플. 744 passed"),
        ("Attack Command Sequencer + PR-2~PR-5 후속 로드맵 (6축 MVP)", "04/17~20", "100%",
            "engage→verify→settle 3-phase 재구성, AttackOutcome Enum. PR-2 Perception·PR-3 Lock-on·PR-4 Stuck Detection·PR-5 Lead + Skill Rotation. 753→941 passed"),
        ("firmware/v2 정합 drift 4건 수정 + PR-6 Attack Strategy 재설계 + PR-7 HP Probe 보정", "04/20", "100%",
            "D1 proto bool attack=8 누락, D2 drag edge-trigger. AttackStrategy Protocol + 3 구현체, default space 키. HPProbeConfig sampling_policy 4종·auto_disable_frames=30. 941→1010 passed"),
    ]),
    ("에이전트·하네스 설계", "자동화 인프라 구축", [
        ("주간업무보고 PPT 생성 에이전트 하네스 설계", "04/16", "100%",
            "레퍼런스 PPT 19슬라이드 구조 분석(맑은 고딕+나눔스퀘어 네오 Heavy, Office 강조색). 단일 에이전트(weekly-report-ppt-generator) + 전용 스킬 분리, reference.pptx + blocks.pptx 듀얼 템플릿"),
        ("일일업무보고 에이전트 프로젝트 그룹핑 고도화 플랜 수립", "04/16", "100%",
            "현재 4개 고정 카테고리 구조를 프로젝트 H2/카테고리 H3 구조로 전환. 프로젝트 식별 7단계 체인, 구 포맷 Preserve 전략, E2E 검증 6종 시나리오"),
    ]),
]

RISKS = [
    ("lineauto pre-existing lint/mypy 위반 누적", "미해결",
     "PR-8 이후 병합 시 베이스라인 교란 우려", "별도 PR로 분리 정리 예정(non-blocking)"),
    ("ESP32 실하드웨어 스모크 테스트 미완료", "미해결",
     "firmware/v2 edge-trigger drag 실환경 검증 부재, 전략 B→C 전환 판단 불가", "실하드웨어 연결 후 스모크 테스트 진행 필요"),
    ("Docker 환경 Testcontainers Flyway IT 실제 실행 부재", "미해결",
     "V1~V43 전버전 회귀 검증이 Skip 상태로 누적", "CI 파이프라인에서 Docker 환경 확보 후 정기 실행"),
    ("테스트 인프라(Vitest + RTL / Playwright) 미도입 (ATL AXIS)", "미해결",
     "프론트엔드 자동 회귀 검증 부재", "스택 결정 후 단계적 도입 필요"),
]

NEXT_WEEK = [
    ("시험 전/후 Phase 기능 런타임 검증 및 UX 보완", "프로덕션 데이터 기반 Phase 필터 정확도·성능 점검"),
    ("ALE 복합 파형 기능 확장", "구간 반복 기능 사용자 피드백 반영, innerRepeat edge-case 대응"),
    ("lineauto PR-8 후속 로드맵", "skill rotation jitter/pre-delay 추가, MOVE_RANDOM jitter, pre-existing lint 정리 PR"),
    ("Strategy 4종 테스트 추가 및 TestSequenceService 테스트 보강", "SignSwap·VoltFlucRipple·RippleImmunity·EmTestSegmentDetector·ResetRipple 커버리지 갭 해소"),
    ("TestFX Presentation E2E 도입 및 JavaFX 21 module opens 검증", "AleMainController 기반 E2E 시나리오 구축"),
    ("Docker 환경 Testcontainers Flyway IT 실제 실행 및 GitHub Actions CI 구성", "V1~V43 회귀 검증 정기화"),
    ("소부장 R&D 사업계획서 최종 취합 및 제출 준비", "본문 8개 섹션 + 과제 개요 + WBS 취합, 제출 양식 최종 검수"),
    ("ATL AXIS 사업화 영역 기능 고도화", "배지 컬러 시맨틱 토큰 완성도 점검, quotation/create 749줄 분해 리팩터"),
]

STATS = [
    ("분석 기간", "2026-04-14 (화) ~ 2026-04-20 (월)"),
    ("일일보고서", "5 / 5 건 (평일 전체, 토·일 제외)"),
    ("완료 업무", "52 건"),
    ("진행 중 업무", "0 건"),
    ("신규 착수", "52 건"),
    ("다음 주 이월", "8 건"),
    ("ISO 주차", "2026-W16 ~ W17"),
]

SPECIAL_NOTES = [
    ("ATL AXIS 배포 환경 구성 전환 (04/16)",
        "Cloudflare + nginx 리버스 프록시 기반으로 프로덕션 도메인(atl-axis.com) 구성 완료. 후속 SSL/TLS Full(strict) 전환·Origin Certificate 발급·방화벽 설정은 사용자 영역"),
    ("Next.js 16 업그레이드 단행 (04/20)",
        "Next.js 15.x Turbopack의 Tailwind v4 오염 버그로 인한 빌드 실패 해소. --turbopack opt-out + .next 자동 정리 predev 훅으로 이중 안전망 확보"),
    ("하네스 기반 자동화 체계 확장",
        "feature-dev-orchestrator·rams-feature-orchestrator·fullstack-orchestrator 3개 파이프라인 상시 활용. designer→backend+frontend→QA→code-reviewer 순서"),
]

# ------------------------------------------------------------
# Slide builders
# ------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = []

    # 1. Title
    def slide_title(prs):
        s = new_slide(prs)
        add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=ACCENT)
        add_rect(s, Inches(0), Inches(5.6), SLIDE_W, Inches(0.08), fill=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.0),
                    "▶ 주간 업무 보고", size=22, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.2),
                    TITLE, size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.6),
                    PERIOD, size=22, color=RGBColor(0xDD, 0xE8, 0xF5))
        add_textbox(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.5),
                    AUTHOR, size=16, color=RGBColor(0xDD, 0xE8, 0xF5))
        return s
    slides.append(slide_title)

    # 2. Executive summary
    def slide_summary(prs):
        s = new_slide(prs)
        add_header(s, "Executive Summary", PERIOD)
        y = Inches(1.7)
        for i, line in enumerate(EXEC_SUMMARY, 1):
            num_box = add_rect(s, Inches(0.7), y, Inches(0.45), Inches(0.45),
                               fill=ACCENT2)
            add_textbox(s, Inches(0.7), y, Inches(0.45), Inches(0.45),
                        str(i), size=16, bold=True,
                        color=RGBColor(0xFF, 0xFF, 0xFF),
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_textbox(s, Inches(1.3), y - Inches(0.03), Inches(11.5), Inches(0.9),
                        line, size=14, color=TEXT)
            y += Inches(0.95)
        return s
    slides.append(slide_summary)

    # 3~N. Projects
    def slide_project_factory(proj_title, proj_subtitle, items):
        def _f(prs):
            s = new_slide(prs)
            add_header(s, proj_title, f"{proj_subtitle} · 완료 {len(items)}건")
            y = Inches(1.65)
            row_h = min(Inches(5.2 / max(len(items), 1)), Inches(0.95))
            for i, (task, date, progress, detail) in enumerate(items, 1):
                # left colored bar
                add_rect(s, Inches(0.6), y, Inches(0.08), row_h - Inches(0.05), fill=ACCENT2)
                # task title
                add_textbox(s, Inches(0.85), y, Inches(8.8), Inches(0.35),
                            f"{i:02d}. {task}", size=12, bold=True, color=ACCENT)
                # date + progress
                add_textbox(s, Inches(9.7), y, Inches(2.0), Inches(0.35),
                            date, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
                add_textbox(s, Inches(11.75), y, Inches(1.1), Inches(0.35),
                            progress, size=11, bold=True, color=OK, align=PP_ALIGN.RIGHT)
                # detail
                add_textbox(s, Inches(0.85), y + Inches(0.33), Inches(12.0), row_h - Inches(0.35),
                            detail, size=10, color=MUTED)
                y += row_h + Inches(0.05)
            return s
        return _f

    for title, subtitle, items in PROJECTS:
        slides.append(slide_project_factory(title, subtitle, items))

    # Special notes
    def slide_special(prs):
        s = new_slide(prs)
        add_header(s, "주간 특이사항", "구조적 전환 및 인프라 변경")
        y = Inches(1.75)
        for head, body in SPECIAL_NOTES:
            add_rect(s, Inches(0.7), y, Inches(0.12), Inches(1.0), fill=ACCENT)
            add_textbox(s, Inches(0.95), y, Inches(11.8), Inches(0.4),
                        head, size=14, bold=True, color=ACCENT)
            add_textbox(s, Inches(0.95), y + Inches(0.38), Inches(11.8), Inches(0.7),
                        body, size=11, color=TEXT)
            y += Inches(1.25)
        return s
    slides.append(slide_special)

    # Risks
    def slide_risks(prs):
        s = new_slide(prs)
        add_header(s, "리스크 및 이슈", "미해결 항목 및 대응 방안")
        y = Inches(1.75)
        for title_, state, impact, action in RISKS:
            add_rect(s, Inches(0.7), y, Inches(0.12), Inches(1.05), fill=DANGER)
            add_textbox(s, Inches(0.95), y, Inches(9.5), Inches(0.35),
                        title_, size=13, bold=True, color=TEXT)
            # state chip
            add_rect(s, Inches(11.3), y + Inches(0.03), Inches(1.4), Inches(0.3), fill=DANGER)
            add_textbox(s, Inches(11.3), y + Inches(0.03), Inches(1.4), Inches(0.3),
                        state, size=10, bold=True,
                        color=RGBColor(0xFF, 0xFF, 0xFF),
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_textbox(s, Inches(0.95), y + Inches(0.37), Inches(11.8), Inches(0.35),
                        f"영향: {impact}", size=10, color=MUTED)
            add_textbox(s, Inches(0.95), y + Inches(0.70), Inches(11.8), Inches(0.35),
                        f"대응: {action}", size=10, color=OK)
            y += Inches(1.20)
        return s
    slides.append(slide_risks)

    # Next week
    def slide_next(prs):
        s = new_slide(prs)
        add_header(s, "다음 주 계획", "2026-04-21 ~ 2026-04-27 (예정)")
        y = Inches(1.70)
        for i, (task, detail) in enumerate(NEXT_WEEK, 1):
            add_rect(s, Inches(0.7), y, Inches(0.45), Inches(0.45), fill=ACCENT)
            add_textbox(s, Inches(0.7), y, Inches(0.45), Inches(0.45),
                        str(i), size=14, bold=True,
                        color=RGBColor(0xFF, 0xFF, 0xFF),
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_textbox(s, Inches(1.3), y - Inches(0.02), Inches(11.7), Inches(0.32),
                        task, size=12, bold=True, color=TEXT)
            add_textbox(s, Inches(1.3), y + Inches(0.30), Inches(11.7), Inches(0.3),
                        detail, size=10, color=MUTED)
            y += Inches(0.65)
        return s
    slides.append(slide_next)

    # Stats
    def slide_stats(prs):
        s = new_slide(prs)
        add_header(s, "주간 통계", "집계 및 분석 지표")
        # 2-col layout of stats
        y0 = Inches(1.9)
        col_w = Inches(5.9)
        for i, (k, v) in enumerate(STATS):
            col = i % 2
            row = i // 2
            x = Inches(0.7 + col * 6.0)
            y = y0 + Inches(row * 0.75)
            add_rect(s, x, y, col_w, Inches(0.6), fill=BG_SOFT)
            add_rect(s, x, y, Inches(0.12), Inches(0.6), fill=ACCENT2)
            add_textbox(s, x + Inches(0.35), y + Inches(0.07), Inches(2.5), Inches(0.5),
                        k, size=11, color=MUTED)
            add_textbox(s, x + Inches(2.8), y + Inches(0.03), Inches(3.0), Inches(0.5),
                        v, size=14, bold=True, color=ACCENT,
                        anchor=MSO_ANCHOR.MIDDLE)
        # closing line
        add_textbox(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4),
                    "※ 사용자 지정 기간 4/14(화)~4/20(월)은 ISO 주차 W16~W17에 걸쳐 있으며, 주말(4/18·4/19)은 일일업무보고 대상에서 제외됨",
                    size=9, color=MUTED)
        return s
    slides.append(slide_stats)

    # End slide
    def slide_end(prs):
        s = new_slide(prs)
        add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=ACCENT)
        add_textbox(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.0),
                    "감사합니다", size=52, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6),
                    "Q & A", size=20,
                    color=RGBColor(0xDD, 0xE8, 0xF5), align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
                    AUTHOR + " · 2026-04-20", size=13,
                    color=RGBColor(0xDD, 0xE8, 0xF5), align=PP_ALIGN.CENTER)
        return s
    slides.append(slide_end)

    total = len(slides)
    for idx, builder in enumerate(slides, 1):
        s = builder(prs)
        # skip footer for title & end slides
        if idx not in (1, total):
            add_footer(s, idx, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({total} slides)")


if __name__ == "__main__":
    build()
